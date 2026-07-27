"""
Etapas 3, 4 y 5 del pipeline: recolección, extracción/limpieza,
chunking, metadatos e indexación vectorial.

Este script:
1. Recorre el directorio de documentos (DOCUMENTS_DIR).
2. Extrae el texto según el formato (md, pdf, docx, xlsx, csv, json, html).
3. Limpia el texto (espacios duplicados, saltos de línea redundantes).
4. Divide el texto en fragmentos (chunks) con solapamiento.
5. Atribuye metadatos a cada fragmento (archivo, categoría, fecha, versión).
6. Genera embeddings y los guarda en una base de datos vectorial Chroma.
"""

import os
import re
import json
import glob
from datetime import datetime

from dotenv import load_dotenv
from langchain_text_splitters import RecursiveCharacterTextSplitter
from langchain_chroma import Chroma
from langchain_huggingface import HuggingFaceEmbeddings
from langchain_core.documents import Document

import pandas as pd

load_dotenv()

DOCUMENTS_DIR = os.getenv("DOCUMENTS_DIR", "documents")
VECTOR_DB_DIR = os.getenv("VECTOR_DB_DIR", "chroma_db")
CHUNK_SIZE = int(os.getenv("CHUNK_SIZE", 800))
CHUNK_OVERLAP = int(os.getenv("CHUNK_OVERLAP", 120))

# Categoría de negocio inferida por nombre de archivo.
# En un proyecto real, esto vendría de metadatos del sistema de origen (etapa 3.2).
CATEGORY_MAP = {
    "politica_privacidad": "Legal y Compliance",
    "politica_reembolsos": "Operacional",
    "preguntas_frecuentes": "Atención al Cliente",
    "guia_envios_entregas": "Operacional",
    "terminos_condiciones": "Legal y Compliance",
    "manual_atencion_cliente": "Operacional",
    "politica_proveedores": "Legal y Compliance",
    "catalogo_productos_ventas": "Datos y Sistemas",
    "presentacion_lunapuntos": "Marketing y Comercial",
    "historial_promociones": "Marketing y Comercial",
    "configuracion_api_publica": "Datos y Sistemas",
    "ayuda_garantias_soporte": "Atención al Cliente",
}


def clean_text(text: str) -> str:
    """Etapa 4.2: limpieza de ruidos comunes de extracción."""
    text = re.sub(r"\n{3,}", "\n\n", text)
    text = re.sub(r"[ \t]{2,}", " ", text)
    text = text.strip()
    return text


def extract_markdown(path: str) -> str:
    with open(path, "r", encoding="utf-8") as f:
        raw = f.read()
    # Quitamos marcas de encabezado markdown pero conservamos el texto (etapa 4.1)
    text = re.sub(r"^#{1,6}\s*", "", raw, flags=re.MULTILINE)
    return clean_text(text)


def extract_csv(path: str) -> str:
    df = pd.read_csv(path)
    lines = []
    headers = list(df.columns)
    for _, row in df.iterrows():
        line = "; ".join(f"{h}: {row[h]}" for h in headers)
        lines.append(line)
    return clean_text("\n".join(lines))


def extract_json(path: str) -> str:
    with open(path, "r", encoding="utf-8") as f:
        data = json.load(f)
    return clean_text(json.dumps(data, ensure_ascii=False, indent=2))


def extract_docx(path: str) -> str:
    import docx
    doc = docx.Document(path)
    return clean_text("\n".join(p.text for p in doc.paragraphs if p.text.strip()))


def extract_pdf(path: str) -> str:
    from pypdf import PdfReader
    reader = PdfReader(path)
    text = "\n".join(page.extract_text() or "" for page in reader.pages)
    return clean_text(text)


def extract_xlsx(path: str) -> str:
    sheets = pd.read_excel(path, sheet_name=None)
    lines = []
    for sheet_name, df in sheets.items():
        lines.append(f"Hoja: {sheet_name}")
        headers = list(df.columns)
        for _, row in df.iterrows():
            lines.append("; ".join(f"{h}: {row[h]}" for h in headers))
    return clean_text("\n".join(lines))


def extract_html(path: str) -> str:
    from bs4 import BeautifulSoup
    with open(path, "r", encoding="utf-8") as f:
        soup = BeautifulSoup(f.read(), "html.parser")
    return clean_text(soup.get_text(separator="\n"))


def extract_pptx(path: str) -> str:
    from pptx import Presentation
    prs = Presentation(path)
    lines = []
    for i, slide in enumerate(prs.slides, start=1):
        lines.append(f"--- Diapositiva {i} ---")
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    text = "".join(run.text for run in para.runs)
                    if text.strip():
                        lines.append(text)
            if shape.has_table:
                for row in shape.table.rows:
                    cells = [cell.text for cell in row.cells]
                    lines.append(" | ".join(cells))
        if slide.has_notes_slide and slide.notes_slide.notes_text_frame.text.strip():
            lines.append(f"(Notas: {slide.notes_slide.notes_text_frame.text.strip()})")
    return clean_text("\n".join(lines))


EXTRACTORS = {
    ".md": extract_markdown,
    ".csv": extract_csv,
    ".json": extract_json,
    ".docx": extract_docx,
    ".pdf": extract_pdf,
    ".xlsx": extract_xlsx,
    ".html": extract_html,
    ".htm": extract_html,
    ".pptx": extract_pptx,
}


def load_documents() -> list[Document]:
    """Etapa 3 y 4.1: recolección y extracción por formato."""
    docs = []
    paths = glob.glob(os.path.join(DOCUMENTS_DIR, "*"))

    if not paths:
        raise FileNotFoundError(
            f"No se encontraron documentos en '{DOCUMENTS_DIR}'. "
            "Agrega al menos un archivo antes de ejecutar la ingesta."
        )

    for path in paths:
        ext = os.path.splitext(path)[1].lower()
        extractor = EXTRACTORS.get(ext)
        if not extractor:
            print(f"[WARN] Formato no soportado, se omite: {path}")
            continue

        filename = os.path.basename(path)
        stem = os.path.splitext(filename)[0]
        category = CATEGORY_MAP.get(stem, "General")

        try:
            text = extractor(path)
        except Exception as e:
            print(f"[ERROR] No se pudo procesar {path}: {e}")
            continue

        if not text.strip():
            print(f"[WARN] Documento vacío tras extracción: {path}")
            continue

        docs.append(
            Document(
                page_content=text,
                metadata={
                    "source": filename,
                    "categoria": category,
                    "fecha_ingesta": datetime.now().isoformat(),
                    "formato": ext.replace(".", ""),
                },
            )
        )
        print(f"[OK] Procesado: {filename} ({category}, {len(text)} caracteres)")

    return docs


def chunk_documents(docs: list[Document]) -> list[Document]:
    """Etapa 4.3: división en fragmentos con solapamiento."""
    splitter = RecursiveCharacterTextSplitter(
        chunk_size=CHUNK_SIZE,
        chunk_overlap=CHUNK_OVERLAP,
        separators=["\n## ", "\n\n", "\n", ". ", " ", ""],
    )
    chunks = splitter.split_documents(docs)

    # Etapa 4.4: metadatos adicionales por fragmento
    for i, chunk in enumerate(chunks):
        chunk.metadata["chunk_id"] = i
        chunk.metadata["longitud"] = len(chunk.page_content)

    return chunks


def build_index(chunks: list[Document]):
    """Etapa 5: generación de embeddings e indexación vectorial.

    Construye el índice desde cero. Solo debe usarse en la primera carga
    (cuando VECTOR_DB_DIR aún no existe). Para altas/bajas de documentos
    individuales tras el arranque, usar index_document()/remove_document(),
    que no requieren borrar el directorio completo (evita bloqueos de
    archivo en Windows, donde SQLite no permite eliminar un .sqlite3 que
    otro proceso -el agente ya cargado- todavía tiene abierto)."""
    print(f"\nGenerando embeddings para {len(chunks)} fragmentos...")
    embeddings = HuggingFaceEmbeddings(
        model_name="sentence-transformers/all-MiniLM-L6-v2",
        # Fuerza CPU explícitamente. Sin esto, en entornos CPU-only como
        # Streamlit Community Cloud, sentence-transformers/torch puede
        # intentar inicializar el modelo en un "meta device" y fallar con
        # NotImplementedError al moverlo ("Cannot copy out of meta tensor").
        model_kwargs={"device": "cpu"},
    )

    vectordb = Chroma.from_documents(
        documents=chunks,
        embedding=embeddings,
        persist_directory=VECTOR_DB_DIR,
        collection_name="lunashop_docs",
        # all-MiniLM-L6-v2 esta optimizado para similitud coseno, no L2
        # (la metrica por defecto de Chroma). Sin esto, similarity_search_
        # with_relevance_scores() devuelve valores mal calibrados y el
        # umbral de confianza descarta contenido relevante.
        collection_metadata={"hnsw:space": "cosine"},
    )
    print(f"Índice vectorial guardado en '{VECTOR_DB_DIR}/'")
    return vectordb


def _get_embeddings():
    return HuggingFaceEmbeddings(
        model_name="sentence-transformers/all-MiniLM-L6-v2",
        model_kwargs={"device": "cpu"},
    )


def _open_vectordb():
    return Chroma(
        persist_directory=VECTOR_DB_DIR,
        embedding_function=_get_embeddings(),
        collection_name="lunashop_docs",
        collection_metadata={"hnsw:space": "cosine"},
    )


def remove_document(filename: str):
    """Elimina del índice todos los fragmentos que pertenezcan a un archivo,
    identificado por su metadata 'source'. No borra el directorio de Chroma,
    solo los vectores de ese documento. Seguro de llamar aunque el archivo
    no tenga fragmentos indexados (no-op)."""
    vectordb = _open_vectordb()
    vectordb._collection.delete(where={"source": filename})


def index_document(path: str):
    """Extrae, limpia, fragmenta e indexa un único archivo (identificado por
    su ruta completa), sin afectar los fragmentos de otros documentos ya
    indexados. Si el archivo ya tenía fragmentos de una versión anterior
    (re-subida), primero los elimina para no dejar duplicados ni chunks
    obsoletos."""
    ext = os.path.splitext(path)[1].lower()
    extractor = EXTRACTORS.get(ext)
    if not extractor:
        raise ValueError(f"Formato no soportado: {ext}")

    filename = os.path.basename(path)
    stem = os.path.splitext(filename)[0]
    category = CATEGORY_MAP.get(stem, "General")

    text = extractor(path)
    if not text.strip():
        raise ValueError(f"El documento quedó vacío tras la extracción: {filename}")

    doc = Document(
        page_content=text,
        metadata={
            "source": filename,
            "categoria": category,
            "fecha_ingesta": datetime.now().isoformat(),
            "formato": ext.replace(".", ""),
        },
    )
    chunks = chunk_documents([doc])

    # Si es una re-subida del mismo nombre, primero quitamos los chunks viejos.
    remove_document(filename)

    vectordb = _open_vectordb()
    vectordb.add_documents(chunks)
    return len(chunks)


def main():
    print("=== Iniciando pipeline de ingesta: LunaShop Agente ===\n")
    docs = load_documents()
    print(f"\n{len(docs)} documentos cargados.")

    chunks = chunk_documents(docs)
    print(f"{len(chunks)} fragmentos generados (chunk_size={CHUNK_SIZE}, overlap={CHUNK_OVERLAP}).")

    build_index(chunks)
    print("\n=== Ingesta completada con éxito ===")


if __name__ == "__main__":
    main()
